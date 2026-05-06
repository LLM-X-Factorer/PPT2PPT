# -*- coding: utf-8 -*-
"""演讲面（presentation）渲染 · v2

设计哲学（详见 docs/v2-roadmap.md）：
- 6 页骨架：cover + product_keypoint + product_milestone + brand_keypoint + brand_milestone + closing
- 复用 v1 的设计系统、utils、cover/closing 版式
- 新增 slide_keypoint / slide_milestone 两个版式（更大字号、更多留白）

数据从 content_presentation.py 读入。换数据 → 跑 `python redesign_presentation.py` → 输出 output-presentation.pptx。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from redesign import (
    SLIDE_W, SLIDE_H,
    INDIGO_900, INDIGO_700, INDIGO_500, INDIGO_300, INDIGO_100, INDIGO_50,
    BG, WHITE, TEXT_PRI, TEXT_SEC, MUTED, LINE,
    GOLD,
    add_textbox, add_rect, fill_slide_bg,
    slide_cover, slide_closing,
)

import content_presentation as C


# ---------- 演讲面专属版式 ----------

def add_presentation_header(slide, section_tag):
    """演讲面顶部：仅一行 section tag + 一根金色细线。比 v1 header 更克制。"""
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.3),
                section_tag, size=11, color=GOLD, bold=True)
    add_rect(slide, Inches(0.7), Inches(0.85), Inches(0.5), Pt(2.5), fill=GOLD)


def slide_keypoint(prs, data):
    """观点页：大字 Action Title + 3 个 numbered points 横排。

    设计目标：台下扫一眼读到结论 + 三个支撑。每页一个观点。
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, BG)
    add_presentation_header(s, data["section"])

    # 大字 Action Title（比 v1 page header 26pt 大很多）
    add_textbox(s, Inches(0.7), Inches(1.25), Inches(11.9), Inches(1.4),
                data["title"], size=40, color=TEXT_PRI, bold=True,
                line_spacing=1.15)
    # 副标题
    add_textbox(s, Inches(0.7), Inches(2.85), Inches(11.9), Inches(0.45),
                data["subtitle"], size=15, color=INDIGO_500)

    # 三个 numbered points 横排
    points = data["points"]
    n = len(points)
    margin = 0.7
    gap = 0.35
    total_w = 13.333 - margin * 2
    card_w = (total_w - gap * (n - 1)) / n
    y = Inches(4.0)
    h = Inches(2.6)

    for i, p in enumerate(points):
        x = Inches(margin + i * (card_w + gap))
        # 大数字（淡灰）
        add_textbox(s, x, y, Inches(1.5), Inches(1.1),
                    p["num"], size=56, color=INDIGO_100, bold=True,
                    line_spacing=0.95)
        # 顶部金色短线
        add_rect(s, x, y + Inches(1.1), Inches(0.4), Pt(2.5), fill=GOLD)
        # head 加粗
        add_textbox(s, x, y + Inches(1.25), Inches(card_w), Inches(0.5),
                    p["head"], size=20, color=TEXT_PRI, bold=True,
                    line_spacing=1.2)
        # desc
        add_textbox(s, x, y + Inches(1.85), Inches(card_w - 0.2), Inches(h.inches - 1.85),
                    p["desc"], size=12, color=TEXT_SEC, line_spacing=1.5)


def slide_milestone(prs, data):
    """节奏页：大字 Action Title + 3 个 lane 卡片（任务 + 阶段产出）。

    设计目标：让台下看到"做什么 + 怎么交付"，不展开任务名细节。
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, BG)
    add_presentation_header(s, data["section"])

    # 大字 Action Title
    add_textbox(s, Inches(0.7), Inches(1.25), Inches(11.9), Inches(1.4),
                data["title"], size=40, color=TEXT_PRI, bold=True,
                line_spacing=1.15)
    add_textbox(s, Inches(0.7), Inches(2.85), Inches(11.9), Inches(0.45),
                data["subtitle"], size=15, color=INDIGO_500)

    # 三个 lane 横排
    lanes = data["lanes"]
    n = len(lanes)
    margin = 0.7
    gap = 0.3
    total_w = 13.333 - margin * 2
    card_w = (total_w - gap * (n - 1)) / n
    y = Inches(3.85)
    h = Inches(2.85)

    for i, lane in enumerate(lanes):
        x = Inches(margin + i * (card_w + gap))
        # 卡片底
        add_rect(s, x, y, Inches(card_w), h, fill=WHITE, line=LINE)
        # 顶部深色块（lane 名）
        add_rect(s, x, y, Inches(card_w), Inches(0.65), fill=INDIGO_700)
        add_textbox(s, x + Inches(0.3), y + Inches(0.05), Inches(card_w - 0.6), Inches(0.55),
                    lane["name"], size=15, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        # 任务名（卡片正中）
        add_textbox(s, x + Inches(0.3), y + Inches(0.85), Inches(card_w - 0.6), Inches(0.65),
                    lane["task"], size=15, color=TEXT_PRI, bold=True,
                    line_spacing=1.3)
        # 金色分割线
        add_rect(s, x + Inches(0.3), y + Inches(1.65), Inches(0.4), Pt(2), fill=GOLD)
        # 阶段产出（5/6/7 月节奏）
        add_textbox(s, x + Inches(0.3), y + Inches(1.8), Inches(card_w - 0.6), Inches(0.25),
                    "5月 → 6月 → 7月", size=9, color=MUTED, bold=True)
        add_textbox(s, x + Inches(0.3), y + Inches(2.05), Inches(card_w - 0.6), Inches(0.7),
                    lane["outcome"], size=11, color=TEXT_SEC, line_spacing=1.5)


# ---------- 主流程 ----------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 6 页：封面 / 产品现状 / 产品节奏 / 品牌现状 / 品牌节奏 / 结尾
    slide_cover(prs, C.COVER)
    slide_keypoint(prs, C.PRODUCT_KEYPOINT)
    slide_milestone(prs, C.PRODUCT_MILESTONE)
    slide_keypoint(prs, C.BRAND_KEYPOINT)
    slide_milestone(prs, C.BRAND_MILESTONE)
    slide_closing(prs, C.CLOSING)

    out = "output-presentation.pptx"
    prs.save(out)
    print(f"✓ saved: {out}")


if __name__ == "__main__":
    build()
