# -*- coding: utf-8 -*-
"""把内容密集的工作汇报 pptx 重新设计为沉稳商务风的可编辑 pptx

数据从 content.py 读入。换数据 → 跑 `python redesign.py` → 输出 output.pptx。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

import content as C

# ---------- 设计系统 ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INDIGO_900 = RGBColor(0x0F, 0x1B, 0x2E)   # 极深蓝（封面 / 章节）
INDIGO_800 = RGBColor(0x14, 0x22, 0x3A)
INDIGO_700 = RGBColor(0x1F, 0x2A, 0x44)   # 主色
INDIGO_500 = RGBColor(0x33, 0x53, 0x80)
INDIGO_300 = RGBColor(0x9D, 0xB2, 0xCE)
INDIGO_100 = RGBColor(0xE7, 0xEC, 0xF4)
INDIGO_50  = RGBColor(0xF2, 0xF5, 0xFA)

BG         = RGBColor(0xF7, 0xF8, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_PRI   = RGBColor(0x1A, 0x23, 0x32)
TEXT_SEC   = RGBColor(0x4B, 0x55, 0x63)
MUTED      = RGBColor(0x9C, 0xA3, 0xAF)
LINE       = RGBColor(0xE5, 0xE7, 0xEB)
LINE_DARK  = RGBColor(0x33, 0x53, 0x80)

GOLD       = RGBColor(0xC2, 0x8C, 0x4C)   # 暖金 · 沉稳点缀
OK         = RGBColor(0x0E, 0x9F, 0x6E)
WARN       = RGBColor(0xE9, 0x73, 0x16)
RISK       = RGBColor(0xDC, 0x26, 0x26)
DOING      = RGBColor(0x33, 0x53, 0x80)
WAIT       = RGBColor(0xC0, 0xC8, 0xD4)

FONT = "Microsoft YaHei"

# ---------- 通用绘制工具 ----------

def set_run(run, text, *, size=14, color=TEXT_PRI, bold=False, italic=False, font=FONT):
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    # 同时设置 east asian / latin 字体，避免中英混排时部分系统回退异常
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", font)


def add_textbox(slide, x, y, w, h, text="", *, size=14, color=TEXT_PRI,
                bold=False, italic=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    if text:
        run = p.add_run()
        set_run(run, text, size=size, color=color, bold=bold, italic=italic, font=font)
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    return s


def add_line(slide, x1, y1, x2, y2, color=LINE, width=0.75):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def fill_slide_bg(slide, color=BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_page_header(slide, title, subtitle=None, section_tag=None):
    """统一的页面顶部：标题 + 副标题 + 顶部细线"""
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(11), Inches(0.55),
                title, size=26, color=TEXT_PRI, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.95), Inches(11), Inches(0.4),
                    subtitle, size=12, color=MUTED)
    if section_tag:
        add_textbox(slide, Inches(0.6), Inches(0.18), Inches(6), Inches(0.25),
                    section_tag, size=9, color=GOLD, bold=True)
    add_rect(slide, Inches(0.6), Inches(1.45), Inches(0.5), Pt(2.5),
             fill=INDIGO_700)
    add_rect(slide, Inches(1.1), Inches(1.45), Inches(11.633), Pt(0.75),
             fill=LINE)


def add_page_footer(slide, idx, total, section=None):
    add_rect(slide, Inches(0.6), Inches(7.05), Inches(12.133), Pt(0.5), fill=LINE)
    if section:
        add_textbox(slide, Inches(0.6), Inches(7.15), Inches(8), Inches(0.3),
                    section, size=9, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.15), Inches(1.733), Inches(0.3),
                f"{idx:02d} / {total:02d}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------- 页面生成函数 ----------

def slide_cover(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, INDIGO_900)
    # 装饰条
    add_rect(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=GOLD)
    # 网格点装饰
    for i in range(8):
        add_rect(s, Inches(11.6 + (i % 4) * 0.35), Inches(5.8 + (i // 4) * 0.35),
                 Inches(0.06), Inches(0.06), fill=INDIGO_500)
    # 标签
    add_textbox(s, Inches(0.9), Inches(1.1), Inches(8), Inches(0.35),
                data["tag"], size=11, color=GOLD, bold=True)
    # 装饰横线
    add_rect(s, Inches(0.9), Inches(1.55), Inches(0.6), Pt(2.5), fill=GOLD)
    # 主标题
    add_textbox(s, Inches(0.9), Inches(1.95), Inches(11), Inches(1.5),
                data["title"], size=58, color=WHITE, bold=True, line_spacing=1.05)
    # 副标题
    add_textbox(s, Inches(0.9), Inches(3.85), Inches(11), Inches(0.6),
                data["subtitle"], size=20, color=INDIGO_300)
    # 底部分隔
    add_rect(s, Inches(0.9), Inches(6.0), Inches(2.0), Pt(1.5), fill=GOLD)
    add_textbox(s, Inches(0.9), Inches(6.15), Inches(8), Inches(0.4),
                data["byline"], size=12, color=INDIGO_300)


def slide_toc(prs, items, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, BG)
    add_page_header(s, "目录", "Contents", section_tag="OVERVIEW")
    add_page_footer(s, idx, total, section="目录")

    # 2x2 网格
    for i, (num, title, sub) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.1)
        y = Inches(2.0 + row * 2.3)
        # 卡片
        add_rect(s, x, y, Inches(5.9), Inches(2.0), fill=WHITE, line=LINE)
        # 大序号
        add_textbox(s, x + Inches(0.4), y + Inches(0.3), Inches(2), Inches(1.2),
                    num, size=64, color=INDIGO_100, bold=True, line_spacing=0.95)
        # 标题
        add_textbox(s, x + Inches(2.0), y + Inches(0.55), Inches(3.7), Inches(0.55),
                    title, size=22, color=TEXT_PRI, bold=True)
        # 副标题
        add_textbox(s, x + Inches(2.0), y + Inches(1.15), Inches(3.7), Inches(0.4),
                    sub, size=12, color=TEXT_SEC)
        # 装饰线
        add_rect(s, x + Inches(2.0), y + Inches(1.0), Inches(0.4), Pt(2),
                 fill=GOLD)


def slide_section(prs, data, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, INDIGO_900)
    # 大数字背景（用纯数字 "01"/"02"，字号控制在不溢出）
    big_num = data["num"].split()[1]
    add_textbox(s, Inches(0.4), Inches(0.8), Inches(11), Inches(6),
                big_num, size=340, color=INDIGO_800, bold=True,
                line_spacing=0.95)
    # 装饰条
    add_rect(s, Inches(0.9), Inches(2.8), Inches(0.6), Pt(2.5), fill=GOLD)
    # 标签
    add_textbox(s, Inches(0.9), Inches(2.45), Inches(6), Inches(0.35),
                data["num"], size=12, color=GOLD, bold=True)
    # 主标题
    add_textbox(s, Inches(0.9), Inches(3.0), Inches(10), Inches(1.5),
                data["title"], size=88, color=WHITE, bold=True, line_spacing=1.0)
    # 副标题
    add_textbox(s, Inches(0.9), Inches(4.6), Inches(10), Inches(0.5),
                data["subtitle"], size=16, color=INDIGO_300, italic=True)


def slide_matrix_product(prs, data, idx, total):
    """产品矩阵：四层级横向卡片"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, BG)
    add_page_header(s, data["title"], data["subtitle"], section_tag="PART 01 · 产品")
    add_page_footer(s, idx, total, section="产品 · 矩阵")

    tiers = data["tiers"]
    n = len(tiers)
    margin = 0.6
    gap = 0.2
    total_w = 13.333 - margin * 2
    card_w = (total_w - gap * (n - 1)) / n

    for i, t in enumerate(tiers):
        x = Inches(margin + i * (card_w + gap))
        y = Inches(1.85)
        h = Inches(5.0)
        # 卡片底
        add_rect(s, x, y, Inches(card_w), h, fill=WHITE, line=LINE)
        # 顶部色条
        bar_color = OK if t["tone"] == "ok" else GOLD
        add_rect(s, x, y, Inches(card_w), Inches(0.12), fill=INDIGO_700)
        # 层级名
        add_textbox(s, x + Inches(0.3), y + Inches(0.35), Inches(card_w - 0.6), Inches(0.45),
                    t["name"], size=22, color=TEXT_PRI, bold=True)
        # 序号点
        dot = add_rect(s, x + Inches(card_w - 0.55), y + Inches(0.4),
                       Inches(0.3), Inches(0.3), fill=INDIGO_50, shape=MSO_SHAPE.OVAL)
        add_textbox(s, x + Inches(card_w - 0.55), y + Inches(0.4), Inches(0.3), Inches(0.3),
                    f"{i+1:02d}", size=10, color=INDIGO_500, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 分割
        add_rect(s, x + Inches(0.3), y + Inches(0.95), Inches(0.4), Pt(2), fill=bar_color)

        # 核心产品标签
        add_textbox(s, x + Inches(0.3), y + Inches(1.2), Inches(card_w - 0.6), Inches(0.25),
                    "核心产品", size=9, color=MUTED, bold=True)
        add_textbox(s, x + Inches(0.3), y + Inches(1.45), Inches(card_w - 0.6), Inches(0.7),
                    t["products"], size=12.5, color=TEXT_PRI, line_spacing=1.4)

        # 客单价
        add_textbox(s, x + Inches(0.3), y + Inches(2.4), Inches(card_w - 0.6), Inches(0.25),
                    "客单价参考", size=9, color=MUTED, bold=True)
        add_textbox(s, x + Inches(0.3), y + Inches(2.65), Inches(card_w - 0.6), Inches(0.5),
                    t["price"], size=13, color=INDIGO_700, bold=True, line_spacing=1.3)

        # 状态/问题
        issue_color = OK if t["tone"] == "ok" else WARN
        add_rect(s, x + Inches(0.3), y + Inches(3.5), Inches(card_w - 0.6), Pt(0.5), fill=LINE)
        add_textbox(s, x + Inches(0.3), y + Inches(3.65), Inches(card_w - 0.6), Inches(0.25),
                    "现状 · 问题", size=9, color=MUTED, bold=True)
        # 状态点
        add_rect(s, x + Inches(0.3), y + Inches(3.95), Inches(0.12), Inches(0.12),
                 fill=issue_color, shape=MSO_SHAPE.OVAL)
        add_textbox(s, x + Inches(0.5), y + Inches(3.92), Inches(card_w - 0.8), Inches(1.0),
                    t["issue"], size=11, color=TEXT_SEC, line_spacing=1.45)


def slide_matrix_brand(prs, data, idx, total):
    """品牌矩阵：四层级竖向带子项"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, BG)
    add_page_header(s, data["title"], data["subtitle"], section_tag="PART 02 · 品牌")
    add_page_footer(s, idx, total, section="品牌 · 矩阵")

    layers = data["layers"]
    n = len(layers)
    margin = 0.6
    gap = 0.2
    total_w = 13.333 - margin * 2
    card_w = (total_w - gap * (n - 1)) / n
    y = Inches(1.85)
    h = Inches(5.0)

    for i, layer in enumerate(layers):
        x = Inches(margin + i * (card_w + gap))
        # 卡片
        add_rect(s, x, y, Inches(card_w), h, fill=WHITE, line=LINE)
        # 顶部深色块
        add_rect(s, x, y, Inches(card_w), Inches(0.85), fill=INDIGO_700)
        # 层级序号
        add_textbox(s, x + Inches(0.3), y + Inches(0.12), Inches(card_w - 0.6), Inches(0.25),
                    f"L{i+1}", size=9, color=GOLD, bold=True)
        # 层级名
        add_textbox(s, x + Inches(0.3), y + Inches(0.35), Inches(card_w - 0.6), Inches(0.45),
                    layer["name"], size=20, color=WHITE, bold=True)

        # 子项
        item_y = y + Inches(0.95)
        item_h = 1.32 if len(layer["items"]) == 3 else 1.85
        for j, item in enumerate(layer["items"]):
            # 序号
            add_textbox(s, x + Inches(0.3), item_y, Inches(0.6), Inches(0.22),
                        f"0{j+1}", size=10, color=GOLD, bold=True)
            # 标题
            add_textbox(s, x + Inches(0.3), item_y + Inches(0.22), Inches(card_w - 0.6), Inches(0.32),
                        item["head"], size=12.5, color=TEXT_PRI, bold=True)
            # 描述
            add_textbox(s, x + Inches(0.3), item_y + Inches(0.55), Inches(card_w - 0.6), Inches(item_h - 0.65),
                        item["desc"], size=10.5, color=TEXT_SEC, line_spacing=1.35)
            # 分割
            if j < len(layer["items"]) - 1:
                add_rect(s, x + Inches(0.3), item_y + Inches(item_h - 0.05),
                         Inches(card_w - 0.6), Pt(0.5), fill=LINE)
            item_y += Inches(item_h)


def slide_problems(prs, data, idx, total, section_tag):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, BG)
    add_page_header(s, data["title"], data["subtitle"], section_tag=section_tag)
    add_page_footer(s, idx, total,
                    section="产品 · 问题" if "产品" in data["title"] else "品牌 · 问题")

    items = data["items"]
    n = len(items)
    margin = 0.6
    gap = 0.3
    total_w = 13.333 - margin * 2
    card_w = (total_w - gap * (n - 1)) / n
    y = Inches(2.2)
    h = Inches(4.4)

    for i, it in enumerate(items):
        x = Inches(margin + i * (card_w + gap))
        # 卡片
        add_rect(s, x, y, Inches(card_w), h, fill=WHITE, line=LINE)
        # 左侧色条
        add_rect(s, x, y, Inches(0.1), h, fill=GOLD)
        # 大数字
        add_textbox(s, x + Inches(0.5), y + Inches(0.4), Inches(card_w - 1), Inches(1.6),
                    it["num"], size=88, color=INDIGO_100, bold=True, line_spacing=1.0)
        # 标题
        add_textbox(s, x + Inches(0.5), y + Inches(1.95), Inches(card_w - 1), Inches(0.6),
                    it["head"], size=22, color=TEXT_PRI, bold=True)
        # 装饰线
        add_rect(s, x + Inches(0.5), y + Inches(2.55), Inches(0.5), Pt(2), fill=INDIGO_700)
        # 描述
        add_textbox(s, x + Inches(0.5), y + Inches(2.85), Inches(card_w - 1), Inches(1.5),
                    it["desc"], size=13, color=TEXT_SEC, line_spacing=1.5)


def slide_okr(prs, data, idx, total, section_tag, section_label):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, BG)
    add_page_header(s, data["title"], data["subtitle"], section_tag=section_tag)
    add_page_footer(s, idx, total, section=section_label)

    rows = data["rows"]
    # 表头
    table_x = Inches(0.6)
    table_y = Inches(1.85)
    col_w = [1.3, 3.0, 5.0, 3.4]  # 层级 / 目标 / KR / 关键动作
    table_w = sum(col_w)

    headers = ["层级", "核心目标 (O)", "关键结果 (KR)", "关键动作"]
    # 表头条
    add_rect(s, table_x, table_y, Inches(table_w), Inches(0.45), fill=INDIGO_700)
    cx = table_x
    for i, h in enumerate(headers):
        add_textbox(s, cx + Inches(0.2), table_y + Inches(0.05),
                    Inches(col_w[i] - 0.4), Inches(0.35),
                    h, size=11, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(col_w[i])

    # 行
    avail_h = 5.0  # 1.85 + 0.45 + 5.0 = 7.30，留 0.2 给页脚
    row_h = avail_h / len(rows)
    ry = table_y + Inches(0.45)
    for r_i, row in enumerate(rows):
        # 行底色（斑马纹）
        if r_i % 2 == 1:
            add_rect(s, table_x, ry, Inches(table_w), Inches(row_h), fill=WHITE)
        else:
            add_rect(s, table_x, ry, Inches(table_w), Inches(row_h), fill=INDIGO_50)

        cx = table_x
        # 层级（带左侧色块）
        add_rect(s, cx, ry + Inches(0.15), Inches(0.06), Inches(row_h - 0.3), fill=GOLD)
        add_textbox(s, cx + Inches(0.2), ry + Inches(0.2),
                    Inches(col_w[0] - 0.3), Inches(row_h - 0.4),
                    row["tier"], size=13, color=INDIGO_700, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(col_w[0])
        # 目标
        add_textbox(s, cx + Inches(0.2), ry + Inches(0.2),
                    Inches(col_w[1] - 0.4), Inches(row_h - 0.4),
                    row["objective"], size=12, color=TEXT_PRI, bold=True,
                    line_spacing=1.35, anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(col_w[1])
        # KR list
        kr_text = ""
        for k_i, kr in enumerate(row["krs"]):
            kr_text += f"KR{k_i+1}  {kr}"
            if k_i < len(row["krs"]) - 1:
                kr_text += "\n"
        tb = add_textbox(s, cx + Inches(0.2), ry + Inches(0.15),
                         Inches(col_w[2] - 0.4), Inches(row_h - 0.3),
                         "", size=10.5, color=TEXT_SEC, line_spacing=1.5,
                         anchor=MSO_ANCHOR.MIDDLE)
        tf = tb.text_frame
        tf.paragraphs[0].text = ""
        for k_i, kr in enumerate(row["krs"]):
            p = tf.paragraphs[0] if k_i == 0 else tf.add_paragraph()
            p.line_spacing = 1.45
            r1 = p.add_run()
            set_run(r1, f"KR{k_i+1}", size=9.5, color=GOLD, bold=True)
            r2 = p.add_run()
            set_run(r2, f"   {kr}", size=10.5, color=TEXT_PRI)
        cx += Inches(col_w[2])
        # 关键动作
        add_textbox(s, cx + Inches(0.2), ry + Inches(0.2),
                    Inches(col_w[3] - 0.4), Inches(row_h - 0.4),
                    row["actions"], size=10.5, color=TEXT_SEC,
                    line_spacing=1.5, anchor=MSO_ANCHOR.MIDDLE)
        # 行底分割线
        add_rect(s, table_x, ry + Inches(row_h) - Pt(0.3),
                 Inches(table_w), Pt(0.5), fill=LINE)
        ry += Inches(row_h)


def slide_gantt(prs, data, idx, total):
    """真正的甘特图（产品 Q1）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, BG)
    add_page_header(s, data["title"], data["subtitle"], section_tag="PART 01 · 产品")
    add_page_footer(s, idx, total, section="产品 · Q1 进度")

    months = data["months"]
    groups = data["groups"]

    # 布局：左 3" 任务名 / 中 9.6" 时间轴
    left_x = Inches(0.6)
    timeline_x = Inches(3.8)
    timeline_w = Inches(8.9)
    top_y = Inches(1.95)
    month_w = timeline_w / len(months)
    row_h = Inches(0.38)
    legend_h = Inches(0.28)

    # 月份表头
    for i, m in enumerate(months):
        mx = timeline_x + month_w * i
        add_textbox(s, mx, top_y - Inches(0.4), month_w, Inches(0.3),
                    m, size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        # 月份分隔虚线
        add_rect(s, mx + month_w, top_y, Pt(0.5),
                 row_h * (sum(len(g["tasks"]) for g in groups) + len(groups)) + Inches(0.3),
                 fill=LINE)

    # 任务行
    cur_y = top_y
    total_tasks = sum(len(g["tasks"]) for g in groups) + len(groups)
    # 顶部细线
    add_rect(s, left_x, top_y, Inches(12.133), Pt(0.75), fill=INDIGO_700)

    for g in groups:
        # 分组行
        add_rect(s, left_x, cur_y, Inches(12.133), row_h, fill=INDIGO_50)
        add_rect(s, left_x, cur_y + Inches(0.08), Inches(0.08), row_h - Inches(0.16), fill=GOLD)
        add_textbox(s, left_x + Inches(0.2), cur_y, Inches(3.4), row_h,
                    g["name"], size=12, color=INDIGO_700, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        cur_y += row_h

        for t in g["tasks"]:
            # 任务名
            add_textbox(s, left_x + Inches(0.4), cur_y, Inches(3.2), row_h,
                        t["name"], size=11, color=TEXT_PRI,
                        anchor=MSO_ANCHOR.MIDDLE)
            # 色条
            sm, em = t["span"]
            bar_x = timeline_x + month_w * (sm - 1)
            bar_w = month_w * (em - sm + 1)
            color = {"done": OK, "doing": DOING, "wait": WAIT}[t["status"]]
            bar_h = Inches(0.22)
            bar_y = cur_y + (row_h - bar_h) / 2
            add_rect(s, bar_x + Inches(0.06), bar_y,
                     bar_w - Inches(0.12), bar_h,
                     fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            # 状态文字
            label = {"done": "✓ 已完成", "doing": "进行中", "wait": "待启动"}[t["status"]]
            label_x = bar_x + bar_w + Inches(0.1)
            if em < 6:
                add_textbox(s, label_x, cur_y, Inches(1.6), row_h,
                            label, size=9, color=color, bold=True,
                            anchor=MSO_ANCHOR.MIDDLE)
            # 里程碑
            if "milestone" in t:
                ms_text, ms_month = t["milestone"]
                ms_x = timeline_x + month_w * (ms_month - 0.5) - Inches(0.08)
                ms_y = cur_y + (row_h - Inches(0.16)) / 2
                add_rect(s, ms_x, ms_y, Inches(0.16), Inches(0.16),
                         fill=GOLD, shape=MSO_SHAPE.DIAMOND)
                add_textbox(s, ms_x - Inches(0.5), ms_y - Inches(0.3),
                            Inches(1.2), Inches(0.25),
                            ms_text, size=8, color=GOLD, bold=True,
                            align=PP_ALIGN.CENTER)
            cur_y += row_h

    # 图例（放在最后任务行下方，避开页脚 7.05）
    legend_y = cur_y + Inches(0.2)
    legends = [("已完成", OK), ("进行中", DOING), ("待启动", WAIT), ("里程碑", GOLD)]
    lx = left_x
    for txt, col in legends:
        if txt == "里程碑":
            add_rect(s, lx, legend_y + Inches(0.05), Inches(0.12), Inches(0.12),
                     fill=col, shape=MSO_SHAPE.DIAMOND)
        else:
            add_rect(s, lx, legend_y + Inches(0.06), Inches(0.32), Inches(0.1),
                     fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(s, lx + Inches(0.45), legend_y, Inches(1.2), legend_h,
                    txt, size=9.5, color=TEXT_SEC, anchor=MSO_ANCHOR.MIDDLE)
        lx += Inches(1.6)


def slide_plan(prs, data, idx, total, section_tag, section_label):
    """季度计划三阶段"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, BG)
    add_page_header(s, data["title"], data["subtitle"], section_tag=section_tag)
    add_page_footer(s, idx, total, section=section_label)

    phases = data["phases"]
    groups = data["groups"]

    # 列宽
    left_w = 2.6  # 板块/任务
    right_w = 13.333 - 0.6 * 2 - left_w
    phase_w = right_w / len(phases)

    # 表头
    header_y = Inches(1.85)
    header_h = Inches(0.55)
    # 左侧表头
    add_rect(s, Inches(0.6), header_y, Inches(left_w), header_h, fill=INDIGO_700)
    add_textbox(s, Inches(0.8), header_y, Inches(left_w - 0.4), header_h,
                "板块 · 任务", size=11, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)
    # 阶段表头
    for i, p in enumerate(phases):
        px = Inches(0.6 + left_w + i * phase_w)
        add_rect(s, px, header_y, Inches(phase_w), header_h, fill=INDIGO_800)
        add_textbox(s, px + Inches(0.2), header_y + Inches(0.05),
                    Inches(phase_w - 0.4), Inches(0.25),
                    p["name"], size=12, color=GOLD, bold=True)
        add_textbox(s, px + Inches(0.2), header_y + Inches(0.27),
                    Inches(phase_w - 0.4), Inches(0.25),
                    p["tag"], size=9, color=INDIGO_300)

    # 行
    rows_y_start = header_y + header_h
    avail_h = 7.0 - 1.85 - 0.55 - 0.05  # ≈ 4.55
    row_h = avail_h / len(groups)

    for r_i, g in enumerate(groups):
        ry = rows_y_start + Inches(row_h * r_i)
        # 行底（斑马）
        bg_color = WHITE if r_i % 2 == 0 else INDIGO_50
        add_rect(s, Inches(0.6), ry,
                 Inches(13.333 - 1.2), Inches(row_h),
                 fill=bg_color)
        # 左侧板块
        add_rect(s, Inches(0.6), ry + Inches(0.18),
                 Inches(0.06), Inches(row_h - 0.36), fill=GOLD)
        add_textbox(s, Inches(0.8), ry + Inches(0.15),
                    Inches(left_w - 0.4), Inches(0.4),
                    g["name"], size=12, color=INDIGO_700, bold=True)
        add_textbox(s, Inches(0.8), ry + Inches(0.55),
                    Inches(left_w - 0.4), Inches(row_h - 0.7),
                    g["task"], size=10.5, color=TEXT_SEC, line_spacing=1.4)

        # 阶段单元格
        for p_i, ph_text in enumerate(g["phases"]):
            px = Inches(0.6 + left_w + p_i * phase_w)
            add_textbox(s, px + Inches(0.2), ry + Inches(0.2),
                        Inches(phase_w - 0.4), Inches(row_h - 0.4),
                        ph_text, size=10, color=TEXT_PRI, line_spacing=1.5)
        # 行底线
        add_rect(s, Inches(0.6), ry + Inches(row_h) - Pt(0.3),
                 Inches(13.333 - 1.2), Pt(0.5), fill=LINE)

    # 列分隔线
    for i in range(len(phases) + 1):
        cx = Inches(0.6 + left_w + i * phase_w)
        add_rect(s, cx, rows_y_start, Pt(0.5),
                 Inches(row_h * len(groups)), fill=LINE)


def slide_accounts(prs, data, idx, total):
    """账号数据表"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, BG)
    add_page_header(s, data["title"], data["subtitle"], section_tag="PART 02 · 品牌")
    add_page_footer(s, idx, total, section="品牌 · Q1 进度")

    accs = data["accounts"]
    # 三列布局
    margin = 0.6
    gap = 0.25
    total_w = 13.333 - margin * 2
    cols = 3
    card_w = (total_w - gap * (cols - 1)) / cols
    rows_per_col = 3
    card_h = (7.0 - 2.0 - 0.4) / rows_per_col  # ≈ 1.53

    for i, a in enumerate(accs):
        col = i % cols
        row = i // cols
        x = Inches(margin + col * (card_w + gap))
        y = Inches(2.0 + row * card_h)
        h = Inches(card_h - 0.15)
        # 卡片
        add_rect(s, x, y, Inches(card_w), h, fill=WHITE, line=LINE)
        # 状态色条
        status_color = OK if a["status"] == "ok" else WARN
        add_rect(s, x, y, Inches(0.1), h, fill=status_color)
        # 平台标签
        add_textbox(s, x + Inches(0.3), y + Inches(0.15),
                    Inches(card_w - 0.6), Inches(0.25),
                    a["platform"], size=9, color=GOLD, bold=True)
        # 账号名
        add_textbox(s, x + Inches(0.3), y + Inches(0.4),
                    Inches(card_w * 0.62), Inches(0.4),
                    a["name"], size=14, color=TEXT_PRI, bold=True)
        # 状态徽章
        badge_text = "正常" if a["status"] == "ok" else "落后"
        badge_w = 0.55
        bx = x + Inches(card_w - badge_w - 0.25)
        add_rect(s, bx, y + Inches(0.42), Inches(badge_w), Inches(0.28),
                 fill=status_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(s, bx, y + Inches(0.42), Inches(badge_w), Inches(0.28),
                    badge_text, size=9, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 数据
        add_textbox(s, x + Inches(0.3), y + Inches(0.85),
                    Inches(card_w - 0.6), Inches(0.35),
                    a["metric"], size=15, color=INDIGO_700, bold=True)
        # 分割线
        add_rect(s, x + Inches(0.3), y + Inches(1.2),
                 Inches(card_w - 0.6), Pt(0.5), fill=LINE)
        # 问题
        add_textbox(s, x + Inches(0.3), y + Inches(1.3),
                    Inches(card_w - 0.6), Inches(card_h - 1.3),
                    a["issue"], size=10, color=TEXT_SEC, line_spacing=1.4)


def slide_closing(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, INDIGO_900)
    # 装饰条
    add_rect(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=GOLD)
    # 居中
    add_textbox(s, Inches(0.6), Inches(2.2), Inches(12.133), Inches(0.4),
                data["tag"], size=12, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.466), Inches(2.65), Inches(0.4), Pt(2.5), fill=GOLD)
    add_textbox(s, Inches(0.6), Inches(2.95), Inches(12.133), Inches(1.5),
                data["title"], size=66, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER, line_spacing=1.05)
    add_textbox(s, Inches(0.6), Inches(4.6), Inches(12.133), Inches(0.5),
                data["subtitle"], size=16, color=INDIGO_300, italic=True,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(6.0), Inches(12.133), Inches(0.6),
                data["footer"], size=24, color=GOLD, bold=True,
                align=PP_ALIGN.CENTER)


# ---------- 主流程 ----------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 14 页：封面 / 目录 / 章节1 / 矩阵 / 问题 / OKR / 甘特 / Q2 / 章节2 / 矩阵 / 问题 / OKR / 账号 / Q2 / 结尾
    # 实际 = 1+1 + 1+5 + 1+5 + 1 = 15
    total = 15

    slide_cover(prs, C.COVER)                                   # 1
    slide_toc(prs, C.TOC, 2, total)                             # 2
    slide_section(prs, C.SECTION_PRODUCT, 3, total)             # 3
    slide_matrix_product(prs, C.PRODUCT_MATRIX, 4, total)       # 4
    slide_problems(prs, C.PRODUCT_PROBLEMS, 5, total, "PART 01 · 产品")  # 5
    slide_okr(prs, C.PRODUCT_OKR, 6, total, "PART 01 · 产品", "产品 · OKR")  # 6
    slide_gantt(prs, C.PRODUCT_Q1, 7, total)                    # 7
    slide_plan(prs, C.PRODUCT_Q2, 8, total, "PART 01 · 产品", "产品 · Q2 计划")  # 8
    slide_section(prs, C.SECTION_BRAND, 9, total)               # 9
    slide_matrix_brand(prs, C.BRAND_MATRIX, 10, total)          # 10
    slide_problems(prs, C.BRAND_PROBLEMS, 11, total, "PART 02 · 品牌")   # 11
    slide_okr(prs, C.BRAND_OKR, 12, total, "PART 02 · 品牌", "品牌 · OKR")  # 12
    slide_accounts(prs, C.BRAND_Q1, 13, total)                  # 13
    slide_plan(prs, C.BRAND_Q2, 14, total, "PART 02 · 品牌", "品牌 · Q2 计划")  # 14
    slide_closing(prs, C.CLOSING)                               # 15

    out = "output.pptx"
    prs.save(out)
    print(f"✓ saved: {out}")


if __name__ == "__main__":
    build()
